import os
import fitz  # PyMuPDF
from pptx import Presentation
try:
    from PIL import Image
    import pytesseract
except ImportError:
    pass # OCR fallback depends on system packages

def parse_pdf(file_path: str) -> str:
    text_content = []
    try:
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text")
            if not text.strip():
                # fallback OCR
                try:
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                    text = pytesseract.image_to_string(img)
                except Exception as e:
                    print(f"OCR failed for page {page_num}: {e}")
            text_content.append(text)
    except Exception as e:
        print(f"Error parsing PDF: {e}")
    return "\n\n".join(text_content)

def parse_pptx(file_path: str) -> str:
    text_content = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text_content.append("\n".join(slide_text))
    except Exception as e:
        print(f"Error parsing PPTX: {e}")
    return "\n\n".join(text_content)

def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        return parse_pdf(file_path)
    elif ext == '.pptx':
        return parse_pptx(file_path)
    return ""
